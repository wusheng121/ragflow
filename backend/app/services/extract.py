import json
import logging
import re
from pathlib import Path

from collections.abc import Callable

from sqlalchemy.orm import Session

from app.config import get_settings
from app.database import Material, Subject
from app.services.ragflow import RagflowClient, RagflowError

logger = logging.getLogger(__name__)

ProgressCallback = Callable[[int, str], None]


def _report(on_progress: ProgressCallback | None, percent: int, message: str) -> None:
    if on_progress:
        on_progress(min(99, max(0, percent)), message)

TYPE_LABELS = {
    "term": "专业术语",
    "theorem": "专业术语",
    "definition": "专业术语",
    "other": "专业术语",
}

GENERIC_TERMS = {
    "我们", "他们", "这个", "那个", "可以", "进行", "通过", "其中", "因此", "所以",
    "如果", "由于", "方法", "问题", "系统", "研究", "分析", "内容", "部分", "过程",
    "结果", "情况", "方面", "主要", "基本", "相关", "不同", "重要", "一般", "具体",
}


def _clamp_summary(text: str, max_chars: int = 20) -> str:
    text = re.sub(r"\s+", "", (text or "").strip())
    text = text.rstrip("。！？!?；;，,")
    if len(text) <= max_chars:
        return text + ("。" if text and not text.endswith(("。", "！", "？")) else "")
    cut = text[:max_chars]
    for i in range(len(cut) - 1, max(len(cut) - 6, 0), -1):
        if cut[i] in "，、；：":
            cut = cut[:i]
            break
    cut = cut.rstrip("，、；：")
    return cut + "。"


def _clamp_detail(text: str, max_sentences: int = 5) -> str:
    text = (text or "").strip()
    if not text:
        return text
    # Remove source suffix from old cards during normalize
    text = re.sub(r"\n\n（来源：[^）]+）\s*$", "", text)
    text = re.sub(r"\n\n（以上内容摘自资料[^）]+）\s*$", "", text)

    parts = re.split(r"([。！？!?])", text)
    sentences: list[str] = []
    buf = ""
    for part in parts:
        buf += part
        if part in "。！？!?":
            sentence = buf.strip()
            if sentence:
                sentences.append(sentence)
            buf = ""
            if len(sentences) >= max_sentences:
                break
    if buf.strip() and len(sentences) < max_sentences:
        tail = buf.strip()
        if not tail.endswith(("。", "！", "？")):
            tail += "。"
        sentences.append(tail)
    return "".join(sentences[:max_sentences])


def _is_professional_term(concept: str) -> bool:
    concept = concept.strip()
    if not concept or concept in GENERIC_TERMS:
        return False
    if _looks_like_formula(concept):
        return False
    if len(concept) == 1:
        return bool(re.fullmatch(r"[\u4e00-\u9fff]", concept))
    return len(concept) >= 2


def _looks_like_formula(text: str) -> bool:
    s = text.strip()
    if not s:
        return True
    if re.search(r"[=∫∑∏√±×÷≤≥≠≈∞\\$]", s):
        return True
    if re.fullmatch(r"[A-Za-z0-9+\-*/^_{}\[\]().,\s]+", s) and any(c in s for c in "=+-*/^"):
        return True
    return False


def read_material_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix in {".txt", ".md", ".markdown"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")[:12000]
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            parts = []
            for page in reader.pages[:30]:
                parts.append(page.extract_text() or "")
            return "\n".join(parts)[:12000]
        except Exception:
            return ""
    if suffix == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            parts: list[str] = []
            for slide in prs.slides[:50]:
                for shape in slide.shapes:
                    text = getattr(shape, "text", "") or ""
                    if text.strip():
                        parts.append(text.strip())
            return "\n".join(parts)[:12000]
        except Exception:
            return ""
    return ""


def _extract_terms(text: str, limit: int = 15) -> list[str]:
    tokens = re.findall(
        r"[\u4e00-\u9fff]{2,10}|[A-Z][A-Za-z0-9\-]{2,}(?:\s[A-Z][A-Za-z0-9\-]+)?",
        text,
    )
    freq: dict[str, int] = {}
    for t in tokens:
        if not _is_professional_term(t):
            continue
        freq[t] = freq.get(t, 0) + 1
    ranked = sorted(freq.keys(), key=lambda k: (-freq[k], -len(k)))
    return ranked[:limit]


def _summarize_from_context(concept: str, ctx: str, source_name: str) -> str:
    ctx = re.sub(r"\s+", " ", (ctx or "").strip())
    if ctx:
        idx = ctx.find(concept)
        snippet = ctx[idx + len(concept) : idx + len(concept) + 40] if idx >= 0 else ctx[:40]
        snippet = re.sub(r"^[，、：:是为指\s]+", "", snippet)
        if snippet:
            return _clamp_summary(snippet)
    return _clamp_summary(f"资料《{source_name}》中的专业术语")


def _detail_from_context(concept: str, ctx: str, source_name: str, full_text: str) -> str:
    ctx = (ctx or "").strip()
    idx = full_text.find(concept)
    extended = full_text[max(0, idx - 200) : idx + 300].strip() if idx >= 0 else ""
    body = extended or ctx or f"该术语出现在资料《{source_name}》中，请结合原文理解。"
    body = re.sub(r"\s+", " ", body)
    return _clamp_detail(body)


def _local_extract(subject_name: str, materials: list[tuple[str, Path]]) -> list[dict]:
    """Fallback: rule-based extraction from local file text."""
    combined_parts: list[str] = []
    source_name = "资料"
    for name, path in materials:
        text = read_material_text(path)
        if text:
            combined_parts.append(text)
            source_name = name

    full_text = "\n\n".join(combined_parts)
    if not full_text.strip():
        raise RagflowError("无法从资料中读取文本内容")

    cards: list[dict] = []
    terms = _extract_terms(full_text)
    for term in terms:
        if not _is_professional_term(term):
            continue
        idx = full_text.find(term)
        ctx = full_text[max(0, idx - 80) : idx + 120] if idx >= 0 else ""
        cards.append(
            {
                "concept": term,
                "type": "term",
                "summary": _summarize_from_context(term, ctx, source_name),
                "detail": _detail_from_context(term, ctx, source_name, full_text),
                "tags": [],
            }
        )

    return cards[:15]


async def _sync_to_ragflow(
    db: Session,
    subject: Subject,
    materials: list[Material],
    on_progress: ProgressCallback | None = None,
) -> tuple[str, list[str], dict[str, str]]:
    client = RagflowClient()

    _report(on_progress, 12, "连接 RAGFlow 知识库…")
    if not subject.ragflow_dataset_id:
        subject.ragflow_dataset_id = await client.create_dataset(f"review-{subject.name}-{subject.id[:8]}")
        db.commit()

    dataset_id = subject.ragflow_dataset_id
    document_ids: list[str] = []
    source_names: dict[str, str] = {}
    to_parse: list[str] = []
    upload_count = 0

    for material in materials:
        if not material.file_path:
            continue
        path = Path(material.file_path)
        if not path.exists():
            continue

        if material.ragflow_document_id:
            document_ids.append(material.ragflow_document_id)
            source_names[material.ragflow_document_id] = material.name
            if material.status != "parsed":
                to_parse.append(material.ragflow_document_id)
            continue

        _report(on_progress, 15 + min(upload_count * 3, 9), f"上传资料：{material.name}")
        doc_id = await client.upload_document(dataset_id, path, material.name)
        upload_count += 1
        material.ragflow_document_id = doc_id
        material.status = "parsing"
        document_ids.append(doc_id)
        source_names[doc_id] = material.name
        to_parse.append(doc_id)

    db.commit()

    if to_parse:
        _report(on_progress, 22, "提交文档解析任务…")
        await client.parse_documents(dataset_id, to_parse)
        await client.wait_documents_ready(
            dataset_id,
            to_parse,
            on_progress=on_progress,
            progress_base=25,
            progress_span=30,
        )
        for material in materials:
            if material.ragflow_document_id in to_parse:
                material.status = "parsed"
        db.commit()

    _report(on_progress, 55, "文档解析完成")
    return dataset_id, document_ids, source_names


def _normalize_cards(raw_items: list[dict], subject_name: str, source_hint: str) -> list[dict]:
    cards: list[dict] = []
    for item in raw_items:
        concept = (item.get("concept") or item.get("name") or "").strip()
        if not concept or not _is_professional_term(concept):
            continue
        ctype = (item.get("type") or "other").lower()
        if ctype == "formula" or _looks_like_formula(concept):
            continue
        if ctype not in TYPE_LABELS:
            ctype = "other"
        summary = (item.get("summary") or item.get("introduction") or "").strip()
        detail = (item.get("detail") or item.get("description") or "").strip()
        if not summary:
            summary = concept
        if not detail:
            detail = summary
        summary = _clamp_summary(summary)
        detail = _clamp_detail(detail)

        cards.append(
            {
                "concept": concept[:500],
                "summary": summary,
                "detail": detail,
                "tags": [],
            }
        )
    return cards


async def _extract_with_llm(subject_name: str, corpus: str, source_hint: str, count: int = 10) -> list[dict]:
    client = RagflowClient()
    raw = await client.chat_extract_concepts(subject_name, corpus, count)
    return _normalize_cards(raw, subject_name, source_hint)[:count]


def _read_local_corpus(materials: list[Material]) -> str:
    sections: list[str] = []
    for material in materials:
        if not material.file_path:
            continue
        text = read_material_text(Path(material.file_path))
        if text.strip():
            sections.append(f"=== 资料：{material.name} ===\n{text}")
    return "\n\n".join(sections)


async def extract_concepts_from_subject(
    db: Session,
    subject: Subject,
    materials: list[Material],
    on_progress: ProgressCallback | None = None,
    count: int = 10,
) -> list[dict]:
    count = max(1, min(50, count))
    settings = get_settings()
    file_materials = [m for m in materials if m.file_path and Path(m.file_path).exists()]
    if not file_materials:
        raise RagflowError("没有可处理的资料文件")

    _report(on_progress, 5, "读取资料文件…")
    source_hint = "、".join(m.name for m in file_materials[:3])
    local_corpus = _read_local_corpus(file_materials)

    if settings.ragflow_enabled:
        try:
            client = RagflowClient()
            _report(on_progress, 10, "同步资料到 RAGFlow…")
            dataset_id, doc_ids, source_names = await _sync_to_ragflow(
                db, subject, file_materials, on_progress=on_progress
            )
            _report(on_progress, 58, "获取文档文本内容…")
            corpus = await client.collect_corpus(dataset_id, doc_ids, source_names)
            if not corpus.strip():
                corpus = local_corpus
            if not corpus.strip():
                raise RagflowError("RAGFlow 解析后未获取到文本内容")

            _report(on_progress, 65, "AI 分析并抽取专业术语…")
            raw = await client.chat_extract_concepts(subject.name, corpus, count)
            _report(on_progress, 80, "整理知识卡片…")
            cards = _normalize_cards(raw, subject.name, source_hint)[:count]
            if cards:
                _report(on_progress, 90, f"已抽取 {len(cards)} 个术语")
                return cards
            logger.warning("RAGFlow LLM returned no valid cards, trying local corpus with LLM")
        except RagflowError as exc:
            logger.warning("RAGFlow extract failed (%s), trying local corpus with LLM", exc)
        except Exception as exc:
            logger.exception("RAGFlow extraction failed, trying local corpus with LLM")
            logger.warning("RAGFlow error: %s", exc)

        if local_corpus.strip():
            try:
                _report(on_progress, 68, "使用本地资料调用 AI 抽取…")
                cards = await _extract_with_llm(subject.name, local_corpus, source_hint, count)
                if cards:
                    _report(on_progress, 90, f"已抽取 {len(cards)} 个术语")
                    return cards
            except Exception as exc:
                logger.warning("LLM extract with local corpus failed: %s", exc)

    if not local_corpus.strip():
        raise RagflowError("无法从资料中读取文本内容，请上传 TXT/MD/PDF/PPTX 格式文件")

    _report(on_progress, 72, "本地规则抽取术语…")
    cards = _normalize_cards(
        _local_extract(subject.name, [(m.name, Path(m.file_path)) for m in file_materials]),
        subject.name,
        source_hint,
    )[:count]
    _report(on_progress, 90, f"已抽取 {len(cards)} 个术语")
    return cards
