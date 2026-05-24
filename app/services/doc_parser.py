from pathlib import Path
from typing import List
import re

from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document


def extract_text_from_pdf(path: Path) -> str:
    try:
        return pdf_extract_text(str(path))
    except Exception:
        return ""


def _looks_like_garbled_text(text: str) -> bool:
    cjk_chars = [ch for ch in text if "\u4e00" <= ch <= "\u9fff"]
    if len(cjk_chars) < 30:
        return False
    common_chars = set(
        "的一是在不了有和人这中大为上个国我以要他时来用们生到作地于出就分对成会可主发年动同工也能下过子说产种面而方后多定行学法所民得经十三之进着等部度家电力里如水化高自二理起小物现实加量都两体制机当使点从业本去把性好应开它合还因由其些然前外天政四日那社义事平形相全表间样与关各重新线内数正心反明看原又么利比或但质气第向道命此变条只没结解问意建月公无系军很情者最立代想已通并提直题程展五果料象员位入常文总次品式活设及管特件长求老头基资边流路级少图山统接知较将组见计别手角期根论运农指几九区强放决西被干做必战先回则任取据处队南给色光门即保治北造百规热领七海口东导器压志世金增争济阶油思术极交受联什认六共权收证改清网络层传输协议主机路由交换广播地址子网掩码吞吐时延带宽丢包重传拥塞控制窗口校验"
    )
    common_count = sum(ch in common_chars for ch in cjk_chars)
    ratio = common_count / max(len(cjk_chars), 1)
    return ratio < 0.35


def extract_text_from_docx(path: Path) -> str:
    doc = Document(path)
    parts: List[str] = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    return "\n".join(parts)


def extract_text_from_file(path: Path) -> str:
    lower = path.suffix.lower()
    if lower == ".pdf":
        return extract_text_from_pdf(path)
    if lower == ".pptx":
        from pptx import Presentation

        prs = Presentation(path)
        paragraphs: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text
                    if txt:
                        paragraphs.append(txt)
        return "\n".join(paragraphs)
    if lower == ".ppt":
        raise ValueError("旧版 PPT（.ppt）暂不支持，请先另存为 PPTX 后再上传")
    if lower in (".docx", ".doc"):
        if lower == ".doc":
            raise ValueError("旧版 Word（.doc）暂不支持，请先另存为 DOCX 后再上传")
        return extract_text_from_docx(path)
    if lower in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="ignore")
    raise ValueError(f"不支持的文件格式: {lower or '未知'}")

